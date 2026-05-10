import json
import math
import os
import re
import textwrap
import time
from io import BytesIO
from urllib.parse import urlencode
from urllib.request import Request, urlopen

from PIL import Image, ImageDraw, ImageFont, ImageFilter


DEFAULT_PALETTE = {
    "navy": "#13213C",
    "ink": "#1F2937",
    "muted": "#64748B",
    "blue": "#2563EB",
    "teal": "#14B8A6",
    "gold": "#F59E0B",
    "red": "#EF4444",
    "bg": "#F8FAFC",
    "white": "#FFFFFF",
}


def build_document(doc_type, content_path, output_path, asset_dir):
    os.makedirs(asset_dir, exist_ok=True)
    with open(content_path, "r", encoding="utf-8") as f:
        raw = f.read()
    payload = _normalize_payload(raw, doc_type)
    payload["_asset_dir"] = asset_dir
    _resolve_assets(payload)

    if doc_type == "pptx":
        _build_pptx(payload, output_path)
    elif doc_type == "docx":
        _build_docx(payload, output_path)
    elif doc_type == "xlsx":
        _build_xlsx(payload, output_path)
    elif doc_type == "pdf":
        _build_pdf(payload, output_path)
    else:
        raise ValueError("Unsupported document type: %s" % doc_type)

    return json.dumps({
        "ok": True,
        "path": output_path,
        "assets": payload.get("_assets", []),
        "asset_count": len(payload.get("_assets", [])),
    }, ensure_ascii=False)


def _normalize_payload(raw, doc_type):
    try:
        data = json.loads(raw) if raw.strip() else {}
    except Exception:
        data = {"title": "Document", "sections": [{"paragraphs": [raw]}]}

    if doc_type == "pptx" and isinstance(data, list):
        data = {"slides": data}
    if not isinstance(data, dict):
        data = {"title": "Document", "sections": [{"paragraphs": [str(data)]}]}

    data.setdefault("title", "Business Document")
    data.setdefault("theme", {})
    data["theme"] = _theme(data.get("theme", {}))

    if doc_type == "pptx":
        if "slides" not in data:
            data["slides"] = _sections_to_slides(data)
    else:
        if "sections" not in data:
            data["sections"] = _slides_to_sections(data.get("slides", []))
    return data


def _theme(theme):
    if isinstance(theme, str):
        theme = {"name": theme}
    theme = theme if isinstance(theme, dict) else {}
    palette = dict(DEFAULT_PALETTE)
    palette.update(theme.get("palette") or {})
    return {
        "name": theme.get("name", "executive"),
        "palette": palette,
        "font": theme.get("font", "Aptos"),
        "background": theme.get("background", "gradient"),
        "tone": theme.get("tone", "commercial"),
    }


def _sections_to_slides(data):
    slides = [{"title": data.get("title", "Business Document"), "subtitle": data.get("subtitle", "")}]
    for section in data.get("sections", []):
        slides.append({
            "title": section.get("heading", "Section"),
            "body": "\n".join(section.get("paragraphs", [])),
            "bullets": section.get("bullets", []),
            "table": section.get("table"),
            "chart": section.get("chart"),
            "image_query": section.get("image_query"),
        })
    return slides


def _slides_to_sections(slides):
    sections = []
    for slide in slides:
        sections.append({
            "heading": slide.get("title", ""),
            "paragraphs": [p for p in [slide.get("subtitle"), slide.get("body")] if p],
            "bullets": slide.get("bullets", []),
            "table": slide.get("table"),
            "chart": slide.get("chart"),
            "image_query": slide.get("image_query"),
        })
    return sections


def _resolve_assets(payload):
    assets = []
    asset_dir = payload["_asset_dir"]
    queries = []
    for q in payload.get("image_queries", []) or []:
        if q:
            queries.append(str(q))
    for item in payload.get("slides", []) or []:
        if item.get("image_query"):
            queries.append(str(item["image_query"]))
    for item in payload.get("sections", []) or []:
        if item.get("image_query"):
            queries.append(str(item["image_query"]))

    for i, url in enumerate(payload.get("image_urls", []) or []):
        path = _download_image(url, os.path.join(asset_dir, "direct_%02d.jpg" % i))
        if path:
            assets.append({"query": "", "url": url, "path": path, "source": "direct"})

    for i, query in enumerate(dict.fromkeys(queries)):
        found = _search_commons_image(query)
        if found:
            path = _download_image(found["url"], os.path.join(asset_dir, "search_%02d.jpg" % i))
            if path:
                found["path"] = path
                assets.append(found)

    payload["_assets"] = assets[:12]


def _search_commons_image(query):
    try:
        params = {
            "action": "query",
            "generator": "search",
            "gsrsearch": query,
            "gsrnamespace": 6,
            "gsrlimit": 6,
            "prop": "imageinfo",
            "iiprop": "url|mime|extmetadata",
            "format": "json",
        }
        url = "https://commons.wikimedia.org/w/api.php?" + urlencode(params)
        pages = (_http_json(url, timeout=10).get("query") or {}).get("pages") or {}
        for page in pages.values():
            info = (page.get("imageinfo") or [{}])[0]
            url = info.get("url", "")
            mime = info.get("mime", "")
            if url and mime.startswith("image/"):
                return {
                    "query": query,
                    "url": url,
                    "title": page.get("title", ""),
                    "source": "Wikimedia Commons",
                }
    except Exception:
        return None
    return None


def _download_image(url, path):
    try:
        img = Image.open(BytesIO(_http_bytes(url, timeout=18))).convert("RGB")
        img.thumbnail((1800, 1200))
        img.save(path, "JPEG", quality=88)
        return path
    except Exception:
        return None


def _http_json(url, timeout=10):
    try:
        import requests
        r = requests.get(url, timeout=timeout, headers={"User-Agent": "MobileClaw/office-document-builder"})
        r.raise_for_status()
        return r.json()
    except Exception:
        req = Request(url, headers={"User-Agent": "MobileClaw/office-document-builder"})
        with urlopen(req, timeout=timeout) as resp:
            return json.loads(resp.read().decode("utf-8", errors="replace"))


def _http_bytes(url, timeout=18):
    try:
        import requests
        r = requests.get(url, timeout=timeout, headers={"User-Agent": "MobileClaw/office-document-builder"})
        r.raise_for_status()
        return r.content
    except Exception:
        req = Request(url, headers={"User-Agent": "MobileClaw/office-document-builder"})
        with urlopen(req, timeout=timeout) as resp:
            return resp.read()


def _asset_for(payload, item=None, index=0):
    assets = payload.get("_assets", [])
    if item and item.get("image_path") and os.path.exists(item["image_path"]):
        return item["image_path"]
    if item and item.get("image_query"):
        q = str(item["image_query"]).lower()
        for asset in assets:
            if q in str(asset.get("query", "")).lower():
                return asset.get("path")
    if assets:
        return assets[index % len(assets)].get("path")
    return None


def _hex(value):
    value = str(value or "#000000").strip()
    if not value.startswith("#"):
        value = "#" + value
    return value


def _rgb(value):
    value = _hex(value).lstrip("#")
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def _font(size=28, bold=False):
    candidates = [
        "/system/fonts/NotoSansCJK-Regular.ttc",
        "/system/fonts/NotoSans-Regular.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size=size)
        except Exception:
            pass
    return ImageFont.load_default()


def _make_background(path, size, theme, title=""):
    palette = theme["palette"]
    base = Image.new("RGB", size, _rgb(palette["bg"]))
    draw = ImageDraw.Draw(base)
    w, h = size
    left = _rgb(palette["navy"])
    right = _rgb(palette["blue"])
    for x in range(w):
        t = x / max(w - 1, 1)
        color = tuple(int(left[i] * (1 - t) + right[i] * t) for i in range(3))
        draw.line((x, 0, x, h), fill=color)
    accent = Image.new("RGBA", size, (0, 0, 0, 0))
    ad = ImageDraw.Draw(accent)
    ad.ellipse((int(w * .62), -int(h * .25), int(w * 1.15), int(h * .55)), fill=(*_rgb(palette["teal"]), 82))
    ad.ellipse((-int(w * .2), int(h * .62), int(w * .4), int(h * 1.2)), fill=(*_rgb(palette["gold"]), 58))
    accent = accent.filter(ImageFilter.GaussianBlur(24))
    base = Image.alpha_composite(base.convert("RGBA"), accent).convert("RGB")
    if title:
        d = ImageDraw.Draw(base)
        d.text((60, h - 96), title[:80], fill=(255, 255, 255), font=_font(30))
    base.save(path, "JPEG", quality=92)
    return path


def _chart_image(chart, path, size=(1000, 560), theme=None):
    theme = theme or _theme({})
    palette = theme["palette"]
    chart = chart or {}
    ctype = chart.get("type", "bar")
    labels = [str(x) for x in chart.get("labels", [])]
    values = chart.get("values", [])
    if not values and chart.get("series"):
        first = chart["series"][0]
        labels = labels or [str(x) for x in first.get("labels", [])]
        values = first.get("values", [])
    values = [float(v or 0) for v in values]
    if not labels:
        labels = ["A", "B", "C", "D"][:len(values) or 4]
    if not values:
        values = [42, 68, 55, 81][:len(labels)]

    img = Image.new("RGB", size, _rgb(palette["white"]))
    d = ImageDraw.Draw(img)
    w, h = size
    title = str(chart.get("title", ""))
    d.text((42, 24), title, fill=_rgb(palette["ink"]), font=_font(34))
    left, top, right, bottom = 80, 95, w - 48, h - 80
    colors = [_rgb(palette["blue"]), _rgb(palette["teal"]), _rgb(palette["gold"]), _rgb(palette["red"])]
    maxv = max(values) if values else 1

    if ctype == "line":
        points = []
        for i, v in enumerate(values):
            x = left + (right - left) * i / max(len(values) - 1, 1)
            y = bottom - (bottom - top) * v / maxv
            points.append((x, y))
        d.line((left, bottom, right, bottom), fill=(220, 226, 235), width=2)
        if len(points) > 1:
            d.line(points, fill=colors[0], width=5)
        for x, y in points:
            d.ellipse((x - 7, y - 7, x + 7, y + 7), fill=colors[1])
    elif ctype in ("donut", "pie"):
        box = (left + 80, top + 5, left + 420, top + 345)
        total = sum(values) or 1
        start = -90
        for i, v in enumerate(values):
            angle = 360 * v / total
            d.pieslice(box, start, start + angle, fill=colors[i % len(colors)])
            start += angle
        if ctype == "donut":
            d.ellipse((box[0] + 85, box[1] + 85, box[2] - 85, box[3] - 85), fill=_rgb(palette["white"]))
        lx = left + 500
        for i, label in enumerate(labels[:8]):
            y = top + i * 42
            d.rectangle((lx, y + 6, lx + 22, y + 28), fill=colors[i % len(colors)])
            d.text((lx + 36, y), "%s  %.1f" % (label, values[i]), fill=_rgb(palette["ink"]), font=_font(24))
    else:
        gap = 18
        bw = max(24, int((right - left - gap * (len(values) - 1)) / max(len(values), 1)))
        d.line((left, bottom, right, bottom), fill=(220, 226, 235), width=2)
        for i, v in enumerate(values):
            x0 = left + i * (bw + gap)
            x1 = x0 + bw
            y0 = bottom - (bottom - top) * v / maxv
            d.rounded_rectangle((x0, y0, x1, bottom), radius=10, fill=colors[i % len(colors)])
            d.text((x0, bottom + 10), labels[i][:10], fill=_rgb(palette["muted"]), font=_font(18))
            d.text((x0, y0 - 30), _fmt(v), fill=_rgb(palette["ink"]), font=_font(20))
    img.save(path, "PNG")
    return path


def _fmt(v):
    if abs(v) >= 1000:
        return "%.1fk" % (v / 1000.0)
    if float(v).is_integer():
        return str(int(v))
    return "%.1f" % v


def _build_pptx(payload, output_path):
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    theme = payload["theme"]
    palette = theme["palette"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slides = payload.get("slides") or [{"title": payload.get("title", "")}]
    for idx, item in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bg_path = os.path.join(payload["_asset_dir"], "slide_bg_%02d.jpg" % idx)
        _make_background(bg_path, (1600, 900), theme, "")
        slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        title = str(item.get("title", payload.get("title", "")))
        subtitle = str(item.get("subtitle", ""))
        title_box = slide.shapes.add_textbox(Inches(.65), Inches(.55), Inches(7.6), Inches(1.0))
        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(36 if idx else 44)
        p.font.color.rgb = RGBColor(255, 255, 255)
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(.7), Inches(1.55), Inches(7.2), Inches(.55))
            sub.text_frame.text = subtitle
            sub.text_frame.paragraphs[0].font.size = Pt(18)
            sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(230, 238, 250)

        body = item.get("body", "")
        bullets = item.get("bullets", []) or []
        if body or bullets:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.65), Inches(2.15), Inches(5.95), Inches(4.55))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.fill.transparency = 8
            card.line.color.rgb = RGBColor(225, 232, 240)
            tx = slide.shapes.add_textbox(Inches(.95), Inches(2.45), Inches(5.35), Inches(3.95))
            tf = tx.text_frame
            tf.clear()
            lines = bullets if bullets else textwrap.wrap(str(body), 72)
            for i, line in enumerate(lines[:8]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(line)
                p.level = 0
                p.font.size = Pt(19)
                p.font.color.rgb = RGBColor(*_rgb(palette["ink"]))

        image_path = _asset_for(payload, item, idx)
        if image_path:
            slide.shapes.add_picture(image_path, Inches(7.15), Inches(1.5), width=Inches(5.45), height=Inches(3.25))

        chart = item.get("chart")
        if chart:
            labels = [str(x) for x in chart.get("labels", [])]
            values = [float(x or 0) for x in chart.get("values", [])]
            if labels and values:
                chart_data = ChartData()
                chart_data.categories = labels
                chart_data.add_series(str(chart.get("series_name", "Value")), values)
                ctype = str(chart.get("type", "bar"))
                xl_type = XL_CHART_TYPE.LINE_MARKERS if ctype == "line" else XL_CHART_TYPE.DOUGHNUT if ctype == "donut" else XL_CHART_TYPE.COLUMN_CLUSTERED
                graphic = slide.shapes.add_chart(xl_type, Inches(7.15), Inches(4.95), Inches(5.45), Inches(1.85), chart_data).chart
                graphic.has_legend = ctype == "donut"
                if graphic.has_legend and graphic.legend is not None:
                    graphic.legend.position = XL_LEGEND_POSITION.RIGHT
                graphic.chart_title.text_frame.text = str(chart.get("title", ""))
            else:
                chart_path = _chart_image(chart, os.path.join(payload["_asset_dir"], "chart_%02d.png" % idx), theme=theme)
                slide.shapes.add_picture(chart_path, Inches(7.15), Inches(4.65), width=Inches(5.45), height=Inches(2.05))

        table = item.get("table")
        if table:
            _add_ppt_table(slide, table, Inches(.8), Inches(5.0), Inches(5.6), Inches(1.45), RGBColor(*_rgb(palette["blue"])))

    prs.save(output_path)


def _add_ppt_table(slide, rows, x, y, w, h, header_color):
    from pptx.util import Pt
    if not rows:
        return
    rows = rows[:6]
    cols = max(len(r) for r in rows)
    shape = slide.shapes.add_table(len(rows), cols, x, y, w, h)
    table = shape.table
    for r, row in enumerate(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(row[c]) if c < len(row) else ""
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color


def _build_docx(payload, output_path):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    theme = payload["theme"]
    palette = theme["palette"]
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(.7)
    section.bottom_margin = Inches(.7)
    section.left_margin = Inches(.72)
    section.right_margin = Inches(.72)
    styles = doc.styles
    styles["Normal"].font.name = theme["font"]
    styles["Normal"].font.size = Pt(10.5)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(str(payload.get("title", "Business Document")))
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(*_rgb(palette["navy"]))
    if payload.get("subtitle"):
        p = doc.add_paragraph(str(payload["subtitle"]))
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    hero = _asset_for(payload, None, 0)
    if hero:
        doc.add_picture(hero, width=Inches(6.4))

    for idx, sec in enumerate(payload.get("sections", [])):
        if sec.get("heading"):
            doc.add_heading(str(sec["heading"]), level=min(int(sec.get("level", 1)), 3))
        for para in sec.get("paragraphs", []):
            doc.add_paragraph(str(para))
        for item in sec.get("bullets", []):
            doc.add_paragraph(str(item), style="List Bullet")
        if sec.get("chart"):
            chart_path = _chart_image(sec["chart"], os.path.join(payload["_asset_dir"], "doc_chart_%02d.png" % idx), theme=theme)
            doc.add_picture(chart_path, width=Inches(6.2))
        img = _asset_for(payload, sec, idx + 1)
        if img and idx % 2 == 0:
            doc.add_picture(img, width=Inches(5.8))
        if sec.get("table"):
            _add_docx_table(doc, sec["table"])

    doc.save(output_path)


def _add_docx_table(doc, rows):
    if not rows:
        return
    cols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=cols)
    table.style = "Table Grid"
    for r, row in enumerate(rows):
        for c in range(cols):
            table.cell(r, c).text = str(row[c]) if c < len(row) else ""


def _build_xlsx(payload, output_path):
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    theme = payload["theme"]
    palette = theme["palette"]
    wb = Workbook()
    raw_sheets = payload.get("sheets")
    if not raw_sheets:
        raw_sheets = [{"name": "Data", "headers": payload.get("headers", []), "rows": payload.get("rows", []), "charts": payload.get("charts", [])}]
    for idx, sheet_data in enumerate(raw_sheets):
        ws = wb.active if idx == 0 else wb.create_sheet()
        ws.title = str(sheet_data.get("name", "Sheet%d" % (idx + 1)))[:31]
        headers = sheet_data.get("headers", [])
        if headers:
            ws.append(headers)
        for row in sheet_data.get("rows", []):
            ws.append(list(row) if isinstance(row, (list, tuple)) else list(row.values()))
        if ws.max_row == 1 and not headers:
            ws.append(["Metric", "Value"])
            for label, value in zip(payload.get("labels", []), payload.get("values", [])):
                ws.append([label, value])
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor=_hex(palette["blue"]).lstrip("#"))
            cell.alignment = Alignment(horizontal="center")
        for col in ws.columns:
            max_len = max((len(str(cell.value)) for cell in col if cell.value is not None), default=8)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 3, 42)
        chart_specs = sheet_data.get("charts") or payload.get("charts") or []
        if chart_specs and ws.max_row >= 2 and ws.max_column >= 2:
            for i, spec in enumerate(chart_specs[:3]):
                chart = LineChart() if spec.get("type") == "line" else PieChart() if spec.get("type") in ("pie", "donut") else BarChart()
                chart.title = spec.get("title", "Chart")
                data = Reference(ws, min_col=2, min_row=1, max_row=ws.max_row)
                cats = Reference(ws, min_col=1, min_row=2, max_row=ws.max_row)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)
                ws.add_chart(chart, "E%d" % (2 + i * 16))
    wb.save(output_path)


def _build_pdf(payload, output_path):
    theme = payload["theme"]
    palette = theme["palette"]
    pages = []
    page_size = (1240, 1754)

    def new_page():
        img = Image.new("RGB", page_size, _rgb(palette["bg"]))
        d = ImageDraw.Draw(img)
        d.rectangle((0, 0, page_size[0], 120), fill=_rgb(palette["navy"]))
        d.rectangle((0, 116, page_size[0], 124), fill=_rgb(palette["teal"]))
        return img, d, 170

    def draw_wrapped(d, text, x, y, font, fill, width_chars=70, line_gap=10, bullet=None):
        prefix = (bullet + " ") if bullet else ""
        lines = []
        for raw in str(text).splitlines() or [""]:
            wrapped = textwrap.wrap(raw, width_chars) or [""]
            lines.extend(wrapped)
        for i, line in enumerate(lines):
            d.text((x, y), (prefix if i == 0 else "  ") + line, font=font, fill=fill)
            y += font.size + line_gap if hasattr(font, "size") else 30
        return y

    page, draw, y = new_page()
    draw.text((70, 42), str(payload.get("title", "Business Document"))[:70], font=_font(40, True), fill=(255, 255, 255))
    if payload.get("subtitle"):
        y = draw_wrapped(draw, payload["subtitle"], 70, y, _font(24), _rgb(palette["muted"]), 72)
        y += 24
    hero = _asset_for(payload, None, 0)
    if hero:
        try:
            im = Image.open(hero).convert("RGB")
            im.thumbnail((1080, 520))
            page.paste(im, (70, y))
            y += im.height + 30
        except Exception:
            pass

    for idx, sec in enumerate(payload.get("sections", [])):
        if y > 1450:
            pages.append(page)
            page, draw, y = new_page()
        if sec.get("heading"):
            draw.text((70, y), str(sec["heading"])[:80], font=_font(30, True), fill=_rgb(palette["blue"]))
            y += 52
        for para in sec.get("paragraphs", []):
            y = draw_wrapped(draw, para, 82, y, _font(22), _rgb(palette["ink"]), 78)
            y += 12
            if y > 1450:
                pages.append(page)
                page, draw, y = new_page()
        for bullet in sec.get("bullets", []):
            y = draw_wrapped(draw, bullet, 105, y, _font(21), _rgb(palette["ink"]), 74, bullet="-")
            y += 8
        if sec.get("chart"):
            chart_path = _chart_image(sec["chart"], os.path.join(payload["_asset_dir"], "pdf_chart_%02d.png" % idx), size=(1000, 560), theme=theme)
            im = Image.open(chart_path).convert("RGB")
            im.thumbnail((1040, 560))
            if y + im.height > 1580:
                pages.append(page)
                page, draw, y = new_page()
            page.paste(im, (90, y))
            y += im.height + 26
        if sec.get("table"):
            rows = sec["table"][:9]
            col_w = int(1040 / max(max(len(r) for r in rows), 1))
            row_h = 42
            if y + row_h * len(rows) > 1580:
                pages.append(page)
                page, draw, y = new_page()
            for r, row in enumerate(rows):
                for c, cell in enumerate(row):
                    x = 90 + c * col_w
                    fill = _rgb(palette["blue"]) if r == 0 else (255, 255, 255)
                    text_fill = (255, 255, 255) if r == 0 else _rgb(palette["ink"])
                    draw.rectangle((x, y, x + col_w, y + row_h), fill=fill, outline=(203, 213, 225))
                    draw.text((x + 10, y + 10), str(cell)[:24], font=_font(17), fill=text_fill)
                y += row_h
            y += 22

    pages.append(page)
    pages[0].save(output_path, "PDF", resolution=150.0, save_all=True, append_images=pages[1:])
